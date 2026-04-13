#!/usr/bin/env python3
"""
Project Machine Pipeline v3.0
HTML -> PNG -> PPTX via Playwright
Each agent has a distinct visual layout signature.
"""

import asyncio, os, tempfile
from pathlib import Path

BG = '#0D1117'
BG_CARD = '#161B22'
BORDER = '#30363D'
TEXT_PRIMARY = '#F0F6FC'
TEXT_SECONDARY = '#C9D1D9'
TEXT_MUTED = '#8B949E'

AGENTS = {
 'researcher':   {'color': '#1D7AA5', 'name': 'Researcher'},
'analyst':      {'color': '#0D9488', 'name': 'Analyst'},
'writer':       {'color': '#BB1B4E', 'name': 'Writer'},
'designer':     {'color': '#D87B2B', 'name': 'Designer'},
'orchestrator': {'color': '#7C3AED', 'name': 'Orchestrator'},
}

def relay_tracker(current_agent):
    order = ['researcher', 'analyst', 'writer', 'designer']
    parts = []
    found = False
    for a in order:
        ag = AGENTS[a]
        if a == current_agent:
            parts.append(f'<span style="color:{ag["color"]};font-weight:700">{ag["name"]} ⟶</span>')
            found = True
        elif not found:
            parts.append(f'<span style="color:{ag["color"]};opacity:0.5">{ag["name"]} ✓</span>')
        else:
            parts.append(f'<span style="color:{TEXT_MUTED}">{ag["name"]} ○</span>')
    return ' &rarr; '.join(parts)

def badge_html(agent_key, stats='', secondary=None):
    a = AGENTS[agent_key]
    b = f'<span style="display:inline-flex;align-items:center;gap:6px;background:{BG_CARD};border:1.5px solid {a["color"]};border-radius:8px;padding:4px 10px;"><span style="font-size:12px;font-weight:700;color:{a["color"]}">{a["name"]}</span><span style="font-size:10px;color:{TEXT_MUTED}">{stats}</span></span>'
    if secondary:
        s = AGENTS[secondary]
        b += f' <span style="display:inline-flex;align-items:center;gap:6px;background:{BG_CARD};border:1.5px solid {s["color"]};border-radius:8px;padding:4px 10px;margin-right:6px;"><span style="font-size:12px;font-weight:700;color:{s["color"]}">{s["name"]}</span></span>'
    return b

def base_html(content, agent_key, slide_num, total, stats='', secondary=None):
    a = AGENTS[agent_key]
    return f'''<!DOCTYPE html>
<html lang="he">
<head><meta charset="UTF-8">
<style>
@import url('https://fonts.googleapis.com/css2?family=Heebo:wght@300;400;500;700;900&display=swap');
*{{margin:0;padding:0;box-sizing:border-box;}}
body{{width:1333px;height:750px;overflow:hidden;font-family:Calibri,'Heebo',sans-serif;
     background:{BG};color:{TEXT_PRIMARY};position:fixed;top:0;left:0;}}
</style></head>
<body>
<div style="position:absolute;top:0;left:0;right:0;height:4px;display:flex;">
    <div style="flex:1;background:#D87B2B"></div>
    <div style="flex:1;background:#BB1B4E"></div>
    <div style="flex:1;background:#1D7AA5"></div>
</div>
<div style="position:absolute;top:4px;left:0;right:0;height:2px;background:{a["color"]};opacity:0.35;"></div>
{content}
<div style="position:absolute;bottom:10px;left:14px;right:14px;display:flex;justify-content:space-between;align-items:center;">
    <div>{badge_html(agent_key, stats, secondary)}</div>
    <div style="display:flex;flex-direction:column;align-items:flex-end;gap:3px;">
        <div style="font-size:9px;color:{TEXT_MUTED};">{relay_tracker(agent_key)}</div>
        <div style="font-size:9px;color:{TEXT_MUTED};">{slide_num} / {total}</div>
    </div>
</div>
</body></html>'''

# RESEARCHER layout: newspaper — banner title with colored left border, 2 equal columns
def html_researcher(title, points, source='', confidence='HIGH', slide_num=1, total=40):
    color = AGENTS['researcher']['color']
    BG_ALT = '#1C2128'
    conf_color = '#22C55E' if confidence == 'HIGH' else '#EAB308' if confidence == 'MEDIUM' else '#EF4444'
    half = max(1, len(points) // 2)
    def col(pts):
        cells = ''
        for i, p in enumerate(pts):
            bg = BG_CARD if i % 2 == 0 else BG_ALT
            cells += f'<div style="flex:1;min-height:80px;background:{bg};border-right:3px solid {color};border-radius:0 6px 6px 0;padding:12px 16px;display:flex;align-items:center;"><div style="font-size:15px;color:{TEXT_SECONDARY};direction:rtl;text-align:right;line-height:1.7;width:100%;">{p}</div></div>'
        return cells
    content = f'''
    <div style="padding:12px 12px 55px;height:calc(100% - 0px);display:grid;grid-template-rows:auto 1fr;gap:0;">
        <div style="border-left:8px solid {color};padding:10px 16px;margin-bottom:10px;display:flex;align-items:center;justify-content:space-between;">
            <div style="display:flex;gap:8px;align-items:center;">
                <span style="font-size:9px;border:1px solid {conf_color};color:{conf_color};padding:2px 6px;border-radius:3px;font-weight:700;">{confidence}</span>
                <span style="font-size:9px;color:{TEXT_MUTED};">{source}</span>
            </div>
            <div style="font-size:28px;font-weight:700;color:{TEXT_PRIMARY};direction:rtl;">{title}</div>
        </div>
        <div style="display:grid;grid-template-columns:1fr 1fr;gap:10px;height:100%;min-height:0;">
            <div style="display:flex;flex-direction:column;gap:6px;height:100%;min-height:0;">{col(points[:half])}</div>
            <div style="display:flex;flex-direction:column;gap:6px;height:100%;min-height:0;">{col(points[half:])}</div>
        </div>
    </div>'''
    return base_html(content, 'researcher', slide_num, total, '30+ חיפושים')

# ANALYST layout: dashboard — large verdict badge, full-width scoring bar, colored table
def html_analyst(title, verdict, rows, score=7.0, slide_num=1, total=40, secondary=None):
    color = AGENTS['analyst']['color']
    BG_ALT = '#1C2128'
    pct = int(score * 10)
    # Header row gets agent color background
    header_row = rows[0] if rows else []
    data_rows = rows[1:] if len(rows) > 1 else []
    header_html = '<tr>' + ''.join(f'<th style="padding:10px 12px;font-size:13px;font-weight:700;color:white;background:{color};direction:rtl;text-align:right;">{cell}</th>' for cell in header_row) + '</tr>' if header_row else ''
    body_html = ''
    for i, row in enumerate(data_rows):
        bg = BG_CARD if i % 2 == 0 else BG_ALT
        body_html += '<tr>' + ''.join(f'<td style="padding:10px 12px;font-size:15px;direction:rtl;text-align:right;color:{TEXT_SECONDARY};background:{bg};">{cell}</td>' for cell in row) + '</tr>'
    content = f'''
    <div style="padding:12px 12px 55px;height:calc(100% - 0px);display:grid;grid-template-rows:auto 16px 1fr;gap:10px;">
        <div style="display:flex;align-items:center;justify-content:space-between;">
            <div style="background:{color};color:white;padding:10px 24px;border-radius:6px;font-size:22px;font-weight:700;">{verdict}</div>
            <div style="font-size:28px;font-weight:700;direction:rtl;">{title}</div>
        </div>
        <div style="display:flex;gap:8px;align-items:center;width:100%;">
            <div style="flex:1;height:16px;background:{BG_CARD};border-radius:8px;overflow:hidden;">
                <div style="width:{pct}%;height:100%;background:linear-gradient(90deg,{color},{color}dd);border-radius:8px;"></div>
            </div>
            <span style="font-size:12px;font-weight:700;color:{TEXT_MUTED};">{score} / 10</span>
        </div>
        <div style="overflow:hidden;border:1px solid {BORDER};border-radius:8px;height:100%;min-height:0;">
            <table style="width:100%;border-collapse:collapse;height:100%;"><thead>{header_html}</thead><tbody>{body_html}</tbody></table>
        </div>
    </div>'''
    return base_html(content, 'analyst', slide_num, total, 'Weighted Scoring', secondary)

# WRITER layout: editorial — top 40% giant quote, bottom 60% flowing paragraph
def html_writer(headline, body, quote='', slide_num=1, total=40):
    color = AGENTS['writer']['color']
    quote_section = f'''
        <div style="flex:4;display:flex;flex-direction:column;justify-content:center;min-height:0;">
            <div style="font-size:13px;font-weight:700;text-transform:uppercase;letter-spacing:2px;color:{color};direction:rtl;text-align:right;margin-bottom:12px;">{headline}</div>
            <div style="border-right:6px solid {color};padding:18px 24px;background:rgba(187,27,78,0.06);border-radius:0 10px 10px 0;">
                <div style="font-size:28px;font-weight:700;line-height:1.4;direction:rtl;text-align:right;color:{TEXT_PRIMARY};">"{quote}"</div>
            </div>
        </div>''' if quote else f'''
        <div style="flex:4;display:flex;align-items:center;min-height:0;">
            <div style="font-size:36px;font-weight:900;line-height:1.2;direction:rtl;text-align:right;color:{TEXT_PRIMARY};border-right:6px solid {color};padding-right:20px;">{headline}</div>
        </div>'''
    content = f'''
    <div style="padding:12px 12px 55px;height:calc(100% - 0px);display:flex;flex-direction:column;gap:0;">
        {quote_section}
        <div style="flex:6;display:flex;align-items:flex-start;padding-top:16px;border-top:1px solid {BORDER};min-height:0;">
            <div style="font-size:17px;line-height:1.9;direction:rtl;text-align:right;color:{TEXT_SECONDARY};overflow:hidden;">{body}</div>
        </div>
    </div>'''
    return base_html(content, 'writer', slide_num, total, '3 זוויות')

# DESIGNER layout: visual-first — left 40% pictogram zone, right 60% 2x2 card grid
def html_designer(title, svg_icon, items, slide_num=1, total=40):
    color = AGENTS['designer']['color']
    items_html = ''.join(f'<div style="background:{BG_CARD};border:1px solid {BORDER};border-top:3px solid {color};border-radius:8px;display:flex;align-items:center;justify-content:center;height:100%;min-height:0;"><div style="font-size:16px;font-weight:700;direction:rtl;color:{TEXT_PRIMARY};">{item}</div></div>' for item in items)
    content = f'''
    <div style="padding:12px 12px 55px;height:calc(100% - 0px);display:grid;grid-template-columns:40% 60%;gap:14px;">
        <div style="display:flex;flex-direction:column;align-items:center;justify-content:center;
                    background:{BG_CARD};border-radius:12px;border:1px solid {BORDER};height:100%;min-height:0;">
            <div style="width:180px;height:180px;display:flex;align-items:center;justify-content:center;">{svg_icon}</div>
            <div style="font-size:28px;font-weight:700;color:{color};margin-top:14px;">{title}</div>
        </div>
        <div style="display:grid;grid-template-columns:1fr 1fr;grid-template-rows:1fr 1fr;gap:10px;height:100%;min-height:0;">
            {items_html}
        </div>
    </div>'''
    return base_html(content, 'designer', slide_num, total, 'Visual-first')

# DIVIDER layout
def html_divider(text, agent_key='researcher', slide_num=1, total=40):
    color = AGENTS[agent_key]['color']
    content = f'''
    <div style="height:100%;display:flex;flex-direction:column;align-items:center;justify-content:center;gap:18px;">
        <div style="font-size:72px;font-weight:900;color:{color};letter-spacing:-2px;direction:rtl;">{text}</div>
        <div style="display:flex;gap:8px;">
            <div style="width:36px;height:3px;background:#D87B2B;border-radius:2px;"></div>
            <div style="width:36px;height:3px;background:#BB1B4E;border-radius:2px;"></div>
            <div style="width:36px;height:3px;background:#1D7AA5;border-radius:2px;"></div>
        </div>
    </div>'''
    return base_html(content, agent_key, slide_num, total)

# STAT CARDS layout (works for any agent)
def html_stats(title, stats, agent_key='researcher', slide_num=1, total=40):
    color = AGENTS[agent_key]['color']
    cards = ''
    for s in stats:
        cards += f'''<div style="background:{BG_CARD};border:1px solid {BORDER};border-top:3px solid {color};
                         border-radius:8px;padding:16px;display:flex;flex-direction:column;align-items:center;justify-content:center;gap:8px;height:100%;">
            <div style="font-size:56px;font-weight:900;color:{color};line-height:1;">{s.get("number","")}</div>
            <div style="font-size:15px;font-weight:700;direction:rtl;text-align:center;">{s.get("label","")}</div>
            <div style="font-size:10px;color:{TEXT_MUTED};direction:rtl;text-align:center;">{s.get("sub","")}</div>
        </div>'''
    content = f'''
    <div style="padding:12px 12px 55px;height:100%;display:grid;grid-template-rows:44px 1fr;gap:12px;">
        <div style="border-bottom:1px solid {BORDER};display:flex;align-items:center;justify-content:flex-end;padding-bottom:8px;">
            <div style="font-size:28px;font-weight:700;direction:rtl;">{title}</div>
        </div>
        <div style="display:grid;grid-template-columns:repeat({len(stats)},1fr);gap:12px;height:100%;">{cards}</div>
    </div>'''
    return base_html(content, agent_key, slide_num, total)

async def build_presentation(slides_html, output_path):
    """Render HTML slides to PNG and assemble PPTX."""
    from playwright.async_api import async_playwright
    from pptx import Presentation
    from pptx.util import Inches

    with tempfile.TemporaryDirectory() as tmpdir:
        # Write HTML files
        for i, html in enumerate(slides_html):
            with open(f"{tmpdir}/slide-{i+1:02d}.html", 'w', encoding='utf-8') as f:
                f.write(html)

        # Render to PNG
        print(f"Rendering {len(slides_html)} slides...")
        png_paths = []
        async with async_playwright() as p:
            browser = await p.chromium.launch(args=['--disable-web-security'])
            for i in range(len(slides_html)):
                html_path = f"{tmpdir}/slide-{i+1:02d}.html"
                png_path = f"{tmpdir}/slide-{i+1:02d}.png"
                page = await browser.new_page()
                await page.set_viewport_size({"width": 1333, "height": 750})
                await page.goto(f"file://{os.path.abspath(html_path)}")
                await page.wait_for_timeout(2800)
                await page.screenshot(path=png_path, clip={"x":0,"y":0,"width":1333,"height":750})
                await page.close()
                png_paths.append(png_path)
                print(f"  OK slide {i+1}")
            await browser.close()

        # Assemble PPTX
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]
        for png in png_paths:
            s = prs.slides.add_slide(blank)
            s.shapes.add_picture(png, 0, 0, prs.slide_width, prs.slide_height)
        prs.save(output_path)
        print(f"Saved: {output_path}")

# EXAMPLE USAGE — replace with actual research content
if __name__ == '__main__':
    import sys
    if len(sys.argv) > 1:
        topic = sys.argv[1]
    else:
        topic = "Project Machine Demo"
    
    output = os.path.expanduser(f"~/Desktop/{topic.replace(' ','-').lower()}.pptx")
    
    # Claude Code will build slides_html list from research content
    # then call: asyncio.run(build_presentation(slides_html, output))
    print(f"Pipeline ready. Build slides_html list and call build_presentation(slides_html, '{output}')")
    print("Available functions: html_researcher, html_analyst, html_writer, html_designer, html_divider, html_stats")
